#!/usr/bin/env python3
"""
slides-ko.py - parchment design system, Korean slide deck generator.

Usage:
  pip install python-pptx --break-system-packages
  python3 slides-ko.py

Output:
  output.pptx (16:9, parchment aesthetic, Nanum Myeongjo serif)

This is a template. Fill in your content and run it directly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════════════════════
# Design system constants
# ═══════════════════════════════════════════════════════════

PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
BRAND_DEEP  = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
CHARCOAL    = RGBColor(0x4d, 0x4c, 0x48)
OLIVE       = RGBColor(0x50, 0x4e, 0x49)
STONE       = RGBColor(0x6b, 0x6a, 0x64)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

# Korean editorial stack. PowerPoint falls back silently if Nanum Myeongjo
# is not installed on the viewing machine — Apple SD Gothic Neo / Malgun Gothic
# / system serif then take over.
SERIF = "NanumMyeongjo"
SANS  = SERIF

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def blank_slide(prs, bg_color=PARCHMENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_card(slide, left, top, width, height,
             fill=IVORY, border=BORDER, border_weight=0.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card


# ═══════════════════════════════════════════════════════════
# Slide templates
# ═══════════════════════════════════════════════════════════

def cover_slide(prs, title, subtitle, author, date):
    s = blank_slide(prs)
    add_text(s, title,
             Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
             font=SERIF, size=48, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.3), Inches(1), weight_pt=1.5)
    add_text(s, subtitle,
             Inches(1), Inches(4.6), Inches(11.33), Inches(0.8),
             font=SANS, size=18, color=OLIVE,
             align=PP_ALIGN.CENTER)
    add_text(s, f"{author} · {date}",
             Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
             font=SANS, size=13, color=STONE,
             align=PP_ALIGN.CENTER)
    return s


def toc_slide(prs, items):
    s = blank_slide(prs)
    add_text(s, "목차",
             Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             font=SERIF, size=34, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(1.8), Inches(11), weight_pt=1)

    for i, item in enumerate(items):
        y = Inches(2.4 + i * 0.9)
        add_text(s, f"0{i+1}",
                 Inches(1.2), y, Inches(1), Inches(0.6),
                 font=SERIF, size=28, color=BRAND)
        add_text(s, item,
                 Inches(2.4), y, Inches(9), Inches(0.6),
                 font=SERIF, size=22, color=NEAR_BLACK,
                 vanchor=MSO_ANCHOR.MIDDLE)
    return s


def chapter_slide(prs, number, title):
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, f"0{number}",
             Inches(0.8), Inches(0.5), Inches(2), Inches(0.8),
             font=SERIF, size=28, color=WHITE)
    add_text(s, title,
             Inches(1), Inches(3), Inches(11.33), Inches(1.5),
             font=SERIF, size=60, color=WHITE,
             align=PP_ALIGN.CENTER)
    return s


def content_slide(prs, eyebrow, title, body, page_num=None):
    s = blank_slide(prs)
    add_text(s, eyebrow.upper(),
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=11, color=STONE)
    add_text(s, title,
             Inches(1.2), Inches(1.2), Inches(11.33), Inches(1.2),
             font=SERIF, size=34, color=NEAR_BLACK)
    add_text(s, body,
             Inches(1.2), Inches(3), Inches(11), Inches(3.5),
             font=SANS, size=18, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def metrics_slide(prs, title, metrics):
    """metrics: [(value, label), ...]"""
    s = blank_slide(prs)
    add_text(s, title,
             Inches(1.2), Inches(0.8), Inches(11), Inches(1),
             font=SERIF, size=30, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(2), Inches(1))

    n = len(metrics)
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total_w = card_w * n + gap * (n - 1)
    start = (SLIDE_W - total_w) / 2

    for i, (value, label) in enumerate(metrics):
        x = start + (card_w + gap) * i
        add_text(s, value,
                 x, Inches(3), card_w, Inches(1.5),
                 font=SERIF, size=56, color=BRAND,
                 align=PP_ALIGN.CENTER)
        add_text(s, label,
                 x, Inches(4.8), card_w, Inches(0.6),
                 font=SANS, size=14, color=OLIVE,
                 align=PP_ALIGN.CENTER)
    return s


def quote_slide(prs, quote, source):
    s = blank_slide(prs)
    add_text(s, f"\u201c{quote}\u201d",
             Inches(1.5), Inches(2.8), Inches(10.33), Inches(2.5),
             font=SERIF, size=30, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER,
             vanchor=MSO_ANCHOR.MIDDLE)
    add_text(s, f" - {source}",
             Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.4),
             font=SANS, size=14, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


def comparison_slide(prs, eyebrow, left_title, left_items, right_title, right_items, page_num=None):
    """Before/After two-column layout. Divider is warm gray. Left column is muted, right is full-weight."""
    s = blank_slide(prs)
    add_text(s, eyebrow.upper(),
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=11, color=STONE)
    divider = s.shapes.add_connector(1,
        Inches(6.67), Inches(1.0),
        Inches(6.67), Inches(6.8))
    divider.line.color.rgb = BORDER
    divider.line.width = Pt(1)
    add_text(s, left_title,
             Inches(1.2), Inches(1.2), Inches(5), Inches(0.8),
             font=SERIF, size=22, color=OLIVE)
    add_text(s, right_title,
             Inches(7.0), Inches(1.2), Inches(5), Inches(0.8),
             font=SERIF, size=22, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.2), Inches(11.5), weight_pt=0.5)
    for i, item in enumerate(left_items[:4]):
        add_text(s, item,
                 Inches(1.2), Inches(2.6 + i * 0.9), Inches(4.9), Inches(0.7),
                 font=SANS, size=17, color=STONE)
    for i, item in enumerate(right_items[:4]):
        add_text(s, item,
                 Inches(7.0), Inches(2.6 + i * 0.9), Inches(5.2), Inches(0.7),
                 font=SANS, size=17, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def pipeline_slide(prs, eyebrow, title, steps, page_num=None):
    """Numbered process steps: 01/02/03 serif numerals + step title + description.
    steps: list of (step_title, step_desc), max 4 steps.
    """
    s = blank_slide(prs)
    add_text(s, eyebrow.upper(),
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=11, color=STONE)
    add_text(s, title,
             Inches(1.2), Inches(1.1), Inches(11), Inches(0.9),
             font=SERIF, size=32, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.15), Inches(11), weight_pt=0.5)

    n = len(steps[:4])
    step_w = Inches(11.5 / n)
    for i, (step_title, step_desc) in enumerate(steps[:4]):
        x = Inches(1.0) + step_w * i
        add_text(s, f"0{i+1}",
                 x, Inches(2.5), step_w, Inches(0.8),
                 font=SERIF, size=42, color=BRAND)
        add_text(s, step_title,
                 x, Inches(3.45), step_w - Inches(0.2), Inches(0.6),
                 font=SERIF, size=19, color=NEAR_BLACK)
        add_text(s, step_desc,
                 x, Inches(4.15), step_w - Inches(0.2), Inches(2.2),
                 font=SANS, size=15, color=OLIVE)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def ending_slide(prs, message, contact):
    s = blank_slide(prs)
    add_text(s, message,
             Inches(1), Inches(3), Inches(11.33), Inches(1.2),
             font=SERIF, size=44, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.5), Inches(1), weight_pt=1.5)
    add_text(s, contact,
             Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
             font=SANS, size=16, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


# ═══════════════════════════════════════════════════════════
# Main - example deck, replace with your content
# ═══════════════════════════════════════════════════════════

def main():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="output.pptx",
                        help="Output PPTX path (default: output.pptx in cwd)")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs,
        title="{{문서 제목}}",
        subtitle="{{한 줄 설명}}",
        author="{{작성자}}",
        date="2026.04")

    toc_slide(prs, items=[
        "{{1장}}",
        "{{2장}}",
        "{{3장}}",
        "Q&A",
    ])

    chapter_slide(prs, 1, "{{장 제목}}")

    content_slide(prs,
        eyebrow="{{장 · 이 페이지}}",
        title="{{한 문장으로 정리한 핵심 주장}}",
        body=("{{18pt 본문 한 단락. 3줄 이내로 유지. "
              "한 슬라이드, 하나의 핵심 메시지. 독자의 주의력이 가장 희소한 자원.}}"),
        page_num=5)

    metrics_slide(prs,
        title="핵심 결과",
        metrics=[
            ("+42%",   "전환율 상승"),
            ("3.8M",   "월간 활성 사용자"),
            ("99.9%",  "가용성 SLA"),
            ("5,000+", "최대 QPS"),
        ])

    quote_slide(prs,
        quote="Good design is as little design as possible.",
        source="Dieter Rams")

    comparison_slide(prs,
        eyebrow="{{장 · 비교}}",
        left_title="{{Before}}",
        left_items=["{{포인트 A}}", "{{포인트 B}}", "{{포인트 C}}"],
        right_title="{{After}}",
        right_items=["{{개선 A}}", "{{개선 B}}", "{{개선 C}}"],
        page_num=8)

    pipeline_slide(prs,
        eyebrow="{{장 · 프로세스}}",
        title="{{프로세스를 단정형으로 정리한 헤드라인}}",
        steps=[
            ("{{1단계}}", "{{1단계 짧은 설명. 2줄 이내.}}"),
            ("{{2단계}}", "{{2단계 짧은 설명. 2줄 이내.}}"),
            ("{{3단계}}", "{{3단계 짧은 설명. 2줄 이내.}}"),
        ],
        page_num=9)

    ending_slide(prs,
        message="감사합니다",
        contact="{{이메일}} · {{웹사이트}}")

    prs.save(args.out)
    print(f"OK: Saved {args.out}")


if __name__ == '__main__':
    main()
